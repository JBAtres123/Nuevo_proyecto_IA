"""
DeepGrader AI — Servicio RAG
Indexación y búsqueda semántica en materiales del docente.

Arquitectura:
  - MySQL  → metadata y estado del material (tabla materiales_curso)
  - SQLite → vectores/embeddings filtrados por curso_id y docente_id

Formatos soportados: PDF, DOCX, PPTX
"""

from __future__ import annotations

import hashlib
import json
import re
import sqlite3
from pathlib import Path
from typing import Optional

import numpy as np
from google import genai
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from config import get_settings

settings = get_settings()

# ── Base de datos SQLite ────────────────────────────────────

DB_PATH = Path("./data/rag.db")
DB_PATH.parent.mkdir(parents=True, exist_ok=True)


def get_db_connection() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    """
    Crea la tabla de embeddings con soporte multi-tenant.
    docente_id y curso_id permiten aislar el contexto de cada curso.
    """
    conn = get_db_connection()

    conn.execute("""
    CREATE TABLE IF NOT EXISTS embeddings (
        id             TEXT PRIMARY KEY,
        docente_id     INTEGER NOT NULL,
        curso_id       INTEGER NOT NULL,
        material_id    INTEGER,
        document       TEXT NOT NULL,
        embedding      TEXT NOT NULL,
        fuente         TEXT,
        chunk_index    INTEGER,
        archivo_path   TEXT
    )
    """)

    # Índices para acelerar búsquedas por curso
    conn.execute("""
    CREATE INDEX IF NOT EXISTS idx_emb_curso
    ON embeddings (docente_id, curso_id)
    """)

    conn.commit()
    conn.close()


init_db()

# ── Cliente Google GenAI ────────────────────────────────────

_genai_client: Optional[genai.Client] = None


def get_genai_client() -> genai.Client:
    global _genai_client
    if _genai_client is None:
        _genai_client = genai.Client(api_key=settings.gemini_api_key)
    return _genai_client


def _get_embedding(text: str) -> list[float]:
    """Genera embedding usando Google GenAI."""
    client = get_genai_client()
    result = client.models.embed_content(
        model=settings.model_embedding,
        contents=text,
    )
    return result.embeddings[0].values


# ── Utilidades vectoriales ──────────────────────────────────

def _normalize_vector(vector: list[float]) -> np.ndarray:
    arr = np.array(vector, dtype=np.float32)
    norm = np.linalg.norm(arr)
    return arr if norm == 0 else arr / norm


def _cosine_similarity(v1: list[float], v2: list[float]) -> float:
    return float(np.dot(_normalize_vector(v1), _normalize_vector(v2)))


def _split_text(
    text: str,
    chunk_size: int = 500,
    overlap: int = 50,
) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    chunks: list[str] = []
    start = 0
    while start < len(text):
        chunks.append(text[start : start + chunk_size])
        start += chunk_size - overlap
    return [c for c in chunks if len(c.strip()) > 20]


# ── Parsers por formato ─────────────────────────────────────

def _parse_pdf(path: str) -> tuple[str, int]:
    """Extrae texto de un PDF. Retorna (texto, num_paginas)."""
    reader = PdfReader(path)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages), len(reader.pages)


def _parse_docx(path: str) -> tuple[str, int]:
    """
    Extrae texto de un DOCX.
    'Páginas' es una estimación: cada 3000 caracteres ≈ 1 página.
    """
    doc = DocxDocument(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Incluir texto de tablas
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text.strip())

    full_text = "\n\n".join(paragraphs)
    estimated_pages = max(1, len(full_text) // 3000)
    return full_text, estimated_pages


def _parse_pptx(path: str) -> tuple[str, int]:
    """
    Extrae texto de un PPTX.
    Cada diapositiva cuenta como una 'página'.
    """
    prs = Presentation(path)
    slides_text: list[str] = []

    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides_text.append("\n".join(parts))

    return "\n\n".join(slides_text), len(prs.slides)


PARSERS = {
    ".pdf":  _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
}

SUPPORTED_EXTENSIONS = set(PARSERS.keys())


def parse_file(file_path: str) -> tuple[str, int]:
    """
    Despacha al parser correcto según la extensión del archivo.
    Lanza ValueError si el formato no está soportado.
    """
    ext = Path(file_path).suffix.lower()
    parser = PARSERS.get(ext)
    if parser is None:
        raise ValueError(
            f"Formato '{ext}' no soportado. "
            f"Usa: {', '.join(SUPPORTED_EXTENSIONS)}"
        )
    return parser(file_path)


# ── Servicio RAG ────────────────────────────────────────────

class RAGService:
    """
    Servicio RAG con aislamiento por docente y curso.

    Flujo:
      1. index_file()  → parsea, chunquea, embebe y guarda en SQLite
      2. search()      → busca en SQLite filtrando por curso_id
      3. delete_material() → borra chunks de un material_id específico
      4. list_sources()    → lista fuentes de un curso
    """

    # ── Indexación ──────────────────────────────────────────

    def index_file(
        self,
        file_path: str,
        docente_id: int,
        curso_id: int,
        material_id: Optional[int] = None,
    ) -> dict:
        """
        Indexa un archivo (PDF / DOCX / PPTX) en SQLite.

        Args:
            file_path:   Ruta al archivo en disco.
            docente_id:  ID del docente propietario.
            curso_id:    ID del curso al que pertenece el material.
            material_id: ID de la fila en materiales_curso (MySQL). Opcional.

        Returns:
            Dict con estadísticas de indexación.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Archivo no encontrado: {file_path}")

        texto, num_paginas = parse_file(file_path)

        if not texto.strip():
            raise ValueError(
                f"No se pudo extraer texto de '{path.name}'. "
                "El archivo puede estar vacío o protegido."
            )

        chunks = _split_text(texto)

        conn = get_db_connection()
        try:
            for i, chunk in enumerate(chunks):
                chunk_id = hashlib.md5(
                    f"{docente_id}:{curso_id}:{file_path}:{i}:{chunk[:50]}".encode()
                ).hexdigest()

                embedding = _get_embedding(chunk)

                conn.execute(
                    """
                    INSERT OR REPLACE INTO embeddings (
                        id, docente_id, curso_id, material_id,
                        document, embedding, fuente, chunk_index, archivo_path
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        chunk_id,
                        docente_id,
                        curso_id,
                        material_id,
                        chunk,
                        json.dumps(embedding),
                        path.name,
                        i,
                        str(path),
                    ),
                )

            conn.commit()
        finally:
            conn.close()

        return {
            "archivo": path.name,
            "formato": path.suffix.lower(),
            "paginas": num_paginas,
            "chunks": len(chunks),
            "backend": "SQLite + NumPy",
        }

    # Alias para compatibilidad con código anterior
    def index_pdf(self, pdf_path: str, docente_id: int = 0, curso_id: int = 0) -> dict:
        return self.index_file(pdf_path, docente_id, curso_id)

    # ── Búsqueda ────────────────────────────────────────────

    def search(
        self,
        query: str,
        curso_id: int,
        docente_id: int,
        k: int = 5,
    ) -> str:
        """
        Busca los chunks más relevantes para una query,
        restringido al curso del docente.

        Args:
            query:      Texto de búsqueda (pregunta del examen).
            curso_id:   Filtro obligatorio de curso.
            docente_id: Filtro de seguridad para el docente.
            k:          Número de resultados a retornar.
        """
        try:
            conn = get_db_connection()
            rows = conn.execute(
                """
                SELECT document, fuente, embedding
                FROM embeddings
                WHERE docente_id = ? AND curso_id = ?
                """,
                (docente_id, curso_id),
            ).fetchall()
            conn.close()

            if not rows:
                return (
                    "⚠️ No hay materiales indexados para este curso. "
                    "El agente calificará sin contexto del curso."
                )

            query_embedding = _get_embedding(query)

            scored = sorted(
                [
                    {
                        "document": row["document"],
                        "fuente": row["fuente"],
                        "similarity": _cosine_similarity(
                            query_embedding, json.loads(row["embedding"])
                        ),
                    }
                    for row in rows
                ],
                key=lambda x: x["similarity"],
                reverse=True,
            )

            parts = [
                f"[Fuente: {r['fuente']} | Relevancia: {round(r['similarity'] * 100, 1)}%]\n{r['document']}"
                for r in scored[:k]
            ]
            return "\n\n---\n\n".join(parts)

        except Exception as exc:
            return f"Error al consultar RAG: {exc}"

    # ── Gestión de fuentes ──────────────────────────────────

    def list_sources(
        self,
        curso_id: Optional[int] = None,
        docente_id: Optional[int] = None,
    ) -> list[str]:
        """Lista archivos indexados. Filtra por curso si se indica."""
        try:
            conn = get_db_connection()
            if curso_id is not None and docente_id is not None:
                rows = conn.execute(
                    """
                    SELECT DISTINCT fuente
                    FROM embeddings
                    WHERE docente_id = ? AND curso_id = ?
                    """,
                    (docente_id, curso_id),
                ).fetchall()
            else:
                rows = conn.execute(
                    "SELECT DISTINCT fuente FROM embeddings"
                ).fetchall()
            conn.close()
            return sorted([r["fuente"] for r in rows if r["fuente"]])
        except Exception:
            return []

    def delete_source(self, filename: str) -> int:
        """Elimina chunks por nombre de archivo (legacy)."""
        try:
            conn = get_db_connection()
            cursor = conn.execute(
                "DELETE FROM embeddings WHERE fuente = ?", (filename,)
            )
            deleted = cursor.rowcount
            conn.commit()
            conn.close()
            return deleted
        except Exception:
            return 0

    def delete_material(
        self,
        material_id: int,
        docente_id: int,
        curso_id: int,
    ) -> int:
        """
        Elimina todos los chunks de un material_id específico.
        Incluye validación de docente_id/curso_id por seguridad.
        """
        try:
            conn = get_db_connection()
            cursor = conn.execute(
                """
                DELETE FROM embeddings
                WHERE material_id = ?
                  AND docente_id = ?
                  AND curso_id   = ?
                """,
                (material_id, docente_id, curso_id),
            )
            deleted = cursor.rowcount
            conn.commit()
            conn.close()
            return deleted
        except Exception:
            return 0

    def get_chunks_count(
        self,
        material_id: int,
        docente_id: int,
        curso_id: int,
    ) -> int:
        """Retorna cuántos chunks tiene un material indexado."""
        try:
            conn = get_db_connection()
            row = conn.execute(
                """
                SELECT COUNT(*) AS cnt
                FROM embeddings
                WHERE material_id = ?
                  AND docente_id  = ?
                  AND curso_id    = ?
                """,
                (material_id, docente_id, curso_id),
            ).fetchone()
            conn.close()
            return row["cnt"] if row else 0
        except Exception:
            return 0


# Singleton global
rag_service = RAGService()